#!/usr/bin/env python3
"""기관 제공 양식(HWP/DOCX/PPTX) 구조 파싱 + 필드 식별.

사용법:
    python parse_template.py <template_path> [--output <output_path>]

출력:
    JSON 형식의 양식 구조 (섹션 목록 + 필드 위치)
"""
import sys
import json
import subprocess
from pathlib import Path


def parse_hwp(path: str) -> dict:
    """HWP 양식을 hwp5txt로 텍스트 추출 후 섹션 구조 파싱."""
    result = subprocess.run(
        ["hwp5txt", path], capture_output=True, text=True, timeout=30
    )
    text = result.stdout
    sections = []
    current_section = None

    for i, line in enumerate(text.split("\n")):
        line = line.strip()
        # 섹션 헤더 감지 패턴
        if any(
            line.startswith(p)
            for p in ["1-", "2-", "3-", "4-", "󰊱", "󰊲", "󰋎", "󰋏", "󰋐", "󰋑", "󰋒"]
        ):
            if current_section:
                sections.append(current_section)
            current_section = {
                "id": f"section-{len(sections)+1}",
                "title": line[:50],
                "line_start": i,
                "content_placeholder": True,
            }
        elif line.startswith("○") or line.startswith("-"):
            if current_section:
                current_section.setdefault("fields", []).append(line)

    if current_section:
        sections.append(current_section)

    return {
        "format": "hwp",
        "path": path,
        "sections": sections,
        "total_sections": len(sections),
    }


def parse_docx(path: str) -> dict:
    """DOCX 양식을 python-docx로 파싱."""
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx"}

    doc = Document(path)
    sections = []

    for i, para in enumerate(doc.paragraphs):
        if para.style.name.startswith("Heading"):
            sections.append(
                {
                    "id": f"section-{len(sections)+1}",
                    "title": para.text,
                    "style": para.style.name,
                    "paragraph_index": i,
                }
            )

    return {
        "format": "docx",
        "path": path,
        "sections": sections,
        "total_sections": len(sections),
    }


def parse_pptx(path: str) -> dict:
    """PPTX 양식을 python-pptx로 슬라이드/플레이스홀더 파싱."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    prs = Presentation(path)
    slides = []

    for i, slide in enumerate(prs.slides):
        placeholders = []
        for ph in slide.placeholders:
            placeholders.append(
                {"idx": ph.placeholder_format.idx, "name": ph.name, "text": ph.text[:100] if ph.text else ""}
            )
        slides.append(
            {
                "slide_number": i + 1,
                "layout": slide.slide_layout.name if slide.slide_layout else "unknown",
                "placeholders": placeholders,
            }
        )

    return {
        "format": "pptx",
        "path": path,
        "slides": slides,
        "total_slides": len(slides),
    }


def detect_and_parse(path: str) -> dict:
    """파일 확장자에 따라 적절한 파서 호출."""
    p = Path(path)
    if p.suffix.lower() == ".hwp":
        return parse_hwp(path)
    elif p.suffix.lower() == ".docx":
        return parse_docx(path)
    elif p.suffix.lower() == ".pptx":
        return parse_pptx(path)
    else:
        return {"error": f"Unsupported format: {p.suffix}"}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python parse_template.py <template_path> [--output <path>]")
        sys.exit(1)

    template_path = sys.argv[1]
    result = detect_and_parse(template_path)

    output_path = None
    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv):
            output_path = sys.argv[idx + 1]

    json_str = json.dumps(result, ensure_ascii=False, indent=2)

    if output_path:
        Path(output_path).write_text(json_str, encoding="utf-8")
        print(f"Saved to {output_path}")
    else:
        print(json_str)
