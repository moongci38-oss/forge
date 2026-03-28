#!/usr/bin/env python3
"""PPT 양식에 슬라이드별 내용 기입 (python-pptx).

사용법:
    python fill_pptx.py <template.pptx> <content_dir> [--output <output.pptx>]
"""
import sys
import json
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def load_slide_mapping(content_dir: str) -> dict:
    """fill-guide/slide-mapping.json 로드."""
    mapping_path = Path(content_dir) / "fill-guide" / "slide-mapping.json"
    if mapping_path.exists():
        return json.loads(mapping_path.read_text(encoding="utf-8"))
    return {"slides": {}}


def fill_placeholder(slide, idx: int, text: str):
    """슬라이드의 특정 플레이스홀더에 텍스트 삽입."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return True
    return False


def fill_presentation(template_path: str, content_dir: str, output_path: str = None):
    """메인 PPT 기입 함수."""
    prs = Presentation(template_path)
    mapping = load_slide_mapping(content_dir)
    filled_count = 0

    for slide_key, slide_info in mapping.get("slides", {}).items():
        slide_num = int(slide_key.replace("slide-", "")) - 1
        if slide_num >= len(prs.slides):
            print(f"  SKIP: {slide_key} — slide not found")
            continue

        slide = prs.slides[slide_num]
        for field in slide_info.get("fields", []):
            ph_idx = field.get("placeholder_idx", 0)
            content_file = Path(field.get("file", ""))
            if content_file.exists():
                content = content_file.read_text(encoding="utf-8")
                content = re.sub(r"^#{1,6}\s+.*$", "", content, flags=re.MULTILINE).strip()
                if fill_placeholder(slide, ph_idx, content[:500]):
                    filled_count += 1
                    print(f"  FILL: {slide_key} ph[{ph_idx}]")

    if not output_path:
        output_path = template_path.replace(".pptx", "-filled.pptx")

    prs.save(output_path)
    print(f"\nSaved: {output_path} ({filled_count} fields filled)")
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python fill_pptx.py <template.pptx> <content_dir> [--output <path>]")
        sys.exit(1)

    template = sys.argv[1]
    content = sys.argv[2]
    output = None
    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv):
            output = sys.argv[idx + 1]

    fill_presentation(template, content, output)
